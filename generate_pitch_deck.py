import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(7, 11, 20)        # #070b14
CARD_BG = RGBColor(15, 23, 42)       # #0f172a
BORDER_COLOR = RGBColor(30, 41, 59)  # #1e293b
CYAN = RGBColor(56, 189, 248)        # #38bdf8
INDIGO = RGBColor(129, 140, 248)     # #818cf8
EMERALD = RGBColor(52, 211, 153)     # #34d399
ROSE = RGBColor(244, 63, 94)         # #f43f5e
WHITE = RGBColor(248, 250, 252)      # #f8fafc
MUTED = RGBColor(148, 163, 184)      # #94a3b8

IMAGES_DIR = r"C:\Users\Hassaan\.gemini\antigravity\brain\c6fea575-e911-40ac-991c-56e4e413dcd4"

def apply_background(slide):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DARK_BG
    bg_shape.line.fill.background()
    return bg_shape

def add_header(slide, tag_text, title_text, subtitle_text):
    # Tag
    tb_tag = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
    tf_tag = tb_tag.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = tag_text.upper()
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = CYAN
    
    # Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.7))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(24)
    p_title.font.bold = True
    p_title.font.color.rgb = WHITE
    
    # Subtitle
    tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = MUTED

blank_slide_layout = prs.slide_layouts[6]

# ==========================================
# SLIDE 1: TITLE / COVER
# ==========================================
s1 = prs.slides.add_slide(blank_slide_layout)
apply_background(s1)

tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_badge = tf1.paragraphs[0]
p_badge.text = "THE WEBMCP CHALLENGE  •  INTERACTIVE MULTI-MODAL TUTOR"
p_badge.font.size = Pt(11)
p_badge.font.bold = True
p_badge.font.color.rgb = CYAN

p_main = tf1.add_paragraph()
p_main.text = "OmniLab MCP: 3D Interactive\nSTEM & Code Learning Studio"
p_main.font.size = Pt(38)
p_main.font.bold = True
p_main.font.color.rgb = WHITE
p_main.space_before = Pt(14)
p_main.space_after = Pt(14)

p_tagline = tf1.add_paragraph()
p_tagline.text = "Transforming chat threads from passive text into live 3D physics playgrounds, Web Audio synthesizers, and algorithmic voxel simulations with deterministic verification."
p_tagline.font.size = Pt(15)
p_tagline.font.color.rgb = MUTED
p_tagline.space_after = Pt(28)

p_meta = tf1.add_paragraph()
p_meta.text = "13 MCP Tools  |  6 Generative UI Widgets  |  Three.js PBR  |  Cannon-es Physics  |  Web Audio API  |  Gemini AI"
p_meta.font.size = Pt(12)
p_meta.font.bold = True
p_meta.font.color.rgb = EMERALD

# ==========================================
# SLIDE 2: THE PROBLEM VS OUR SOLUTION
# ==========================================
s2 = prs.slides.add_slide(blank_slide_layout)
apply_background(s2)
add_header(s2, "Challenge & Paradigm Shift", "Why Conventional AI Assistants Fall Short in STEM & Coding", "Bridging the gap between static text hallucinations and experiential physical understanding.")

# Card 1: The Problem
card1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.5))
card1.fill.solid()
card1.fill.fore_color.rgb = CARD_BG
card1.line.color.rgb = ROSE
tf_c1 = card1.text_frame
tf_c1.word_wrap = True
tf_c1.margin_top = Inches(0.3)
tf_c1.margin_left = Inches(0.4)
tf_c1.margin_right = Inches(0.4)

p = tf_c1.paragraphs[0]
p.text = "❌ Current AI Chat Limitations"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ROSE

items_prob = [
    ("Wall of Static Text:", "Students glaze over long prose without building spatial or mechanical mental models."),
    ("Hallucinated Grading (Sycophancy):", "LLMs often tell students their broken code is correct, reinforcing critical misunderstandings."),
    ("Zero Sensory Feedback:", "Concepts like acoustics, wave harmonics, torque, and Big-O sorting remain abstract symbols on paper."),
    ("Passive Consumption:", "Students read answers instead of experimenting with variables and discovering principles actively.")
]
for title, desc in items_prob:
    p = tf_c1.add_paragraph()
    p.text = f"• {title} {desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.space_before = Pt(10)

# Card 2: The Solution
card2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.6), Inches(4.5))
card2.fill.solid()
card2.fill.fore_color.rgb = CARD_BG
card2.line.color.rgb = CYAN
tf_c2 = card2.text_frame
tf_c2.word_wrap = True
tf_c2.margin_top = Inches(0.3)
tf_c2.margin_left = Inches(0.4)
tf_c2.margin_right = Inches(0.4)

p = tf_c2.paragraphs[0]
p.text = "✓ The OmniLab MCP Paradigm"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = CYAN

items_sol = [
    ("Generative 3D UI in Chat:", "Renders live Three.js PBR viewports, Cannon-es physics, and audio synths directly inside the chat thread."),
    ("Deterministic Trust Boundary:", "Grading is strictly handled by isolated Python subprocesses and math engines—never left to AI hallucinations."),
    ("Socratic Scaffolding:", "Gemini inspects ground-truth unit test diffs to ask targeted guiding questions rather than giving away answers."),
    ("In-Scene Scientific Explanations:", "Every 3D simulation includes live embedded concept cards with exact formulas and principles.")
]
for title, desc in items_sol:
    p = tf_c2.add_paragraph()
    p.text = f"• {title} {desc}"
    p.font.size = Pt(12)
    p.font.color.rgb = MUTED
    p.space_before = Pt(10)

# ==========================================
# SLIDE 3: COMPLETE ARCHITECTURE & TECH STACK
# ==========================================
s3 = prs.slides.add_slide(blank_slide_layout)
apply_background(s3)
add_header(s3, "Engineering & Architecture", "How OmniLab Works: Protocol & Technologies", "A robust decoupling of AI proposal, protocol messaging, and deterministic execution.")

tech_cards = [
    ("MCP Streamable-HTTP Protocol", "Standardized Model Context Protocol transport on port 3000 handling JSON-RPC tool dispatch and UI resource resolution.", CYAN, Inches(0.8), Inches(2.2)),
    ("Ext-Apps Generative UI", "Client-side iframe protocol (@modelcontextprotocol/ext-apps) enabling bidirectional state exchange and zero-install client rendering.", INDIGO, Inches(4.8), Inches(2.2)),
    ("Google Gemini Multi-Model AI", "Leverages gemini-flash-lite-latest with fallback to generate 3D spatial scene specifications and Socratic tutoring hints.", EMERALD, Inches(8.8), Inches(2.2)),
    ("Three.js WebGL (PBR Shaders)", "Physically-Based Rendering with ACESFilmic tone mapping, PCF soft shadows, MeshPhysicalMaterial clearcoat, and procedural maps.", CYAN, Inches(0.8), Inches(4.6)),
    ("Cannon-es 3D Rigid Physics", "60 Hz continuous collision dynamics engine computing linear momentum, impulse, torque, and variable gravity presets.", ROSE, Inches(4.8), Inches(4.6)),
    ("Deterministic Python Sandbox", "Subprocess execution (python -I) with 5s timeout; checks real assertions before progress_store.py awards verified mastery.", EMERALD, Inches(8.8), Inches(4.6))
]

for title, desc, col, x, y in tech_cards:
    c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(2.1))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = BORDER_COLOR
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = col
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(8)

# ==========================================
# SLIDE 4: SYSTEM 1 - 3D RIGID-BODY PHYSICS PLAYGROUND
# ==========================================
s4 = prs.slides.add_slide(blank_slide_layout)
apply_background(s4)
add_header(s4, "Feature Highlight 01", "3D Rigid-Body Physics Playground (Cannon-es + PBR)", "Hands-on kinetic momentum, collision dynamics, and gravitational experimentation.")

tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.8), Inches(4.6))
tf4 = tb4.text_frame
tf4.word_wrap = True

bullets_s4 = [
    ("60 Hz Collision Engine:", "Simulates mass, velocity vectors, friction (μ=0.28), and restitution (e=0.45) across multiple interacting rigid bodies."),
    ("Interactive Cannonball Blaster:", "Users click anywhere in 3D space to fire high-velocity steel projectiles into balanced dominoes and block towers."),
    ("Variable Gravity Presets:", "Instantly switch gravity: 🌍 Earth (-9.82 m/s²), 🌕 Moon (-1.62 m/s²), 🌌 Zero-G (0.0 m/s²), 🪐 Jupiter (-24.79 m/s²)."),
    ("Embedded Physics Law Card:", "Displays exact formulas for linear momentum (p = m·v), torque (τ = r × F), and cascade potential-to-kinetic conversion.")
]
for t, d in bullets_s4:
    p = tf4.add_paragraph() if tf4.paragraphs[0].text else tf4.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

img_phys = os.path.join(IMAGES_DIR, "3d_physics_explained.png")
if os.path.exists(img_phys):
    s4.shapes.add_picture(img_phys, Inches(5.8), Inches(2.2), Inches(6.7), Inches(4.5))

# ==========================================
# SLIDE 5: SYSTEM 2 - WEB AUDIO FREQUENCY SYNTHESIZER
# ==========================================
s5 = prs.slides.add_slide(blank_slide_layout)
apply_background(s5)
add_header(s5, "Feature Highlight 02", "Interactive Web Audio Synthesizer & 60 FPS Oscilloscope", "Bridging acoustics, Fourier series, and sensory auditory feedback.")

tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.8), Inches(4.6))
tf5 = tb5.text_frame
tf5.word_wrap = True

bullets_s5 = [
    ("Native AudioContext Engine:", "Generates real audio frequencies directly in the user's browser without external audio files."),
    ("Live 60 FPS Oscilloscope:", "Extracts time-domain data via AnalyserNode to draw glowing real-time vibrating waveform harmonics."),
    ("Harmonic Waveform Selectors:", "Switch between Sine ∿, Square ⊓, Sawtooth ⧘, and Triangle ⋀ to observe harmonic timbre changes."),
    ("Interactive Piano Keyboard:", "Playable 8-key piano keyboard (C4 through C5) with resonant low-pass filter cutoff and pitch sliders."),
    ("Acoustic Equations:", "Displays Fourier wave expansions and harmonic frequency relationships alongside audible sound.")
]
for t, d in bullets_s5:
    p = tf5.add_paragraph() if tf5.paragraphs[0].text else tf5.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

img_audio = os.path.join(IMAGES_DIR, "dashboard_audio.png")
if os.path.exists(img_audio):
    s5.shapes.add_picture(img_audio, Inches(5.8), Inches(2.2), Inches(6.7), Inches(4.5))

# ==========================================
# SLIDE 6: SYSTEM 3 - 3D ALGORITHM VOXEL VISUALIZER
# ==========================================
s6 = prs.slides.add_slide(blank_slide_layout)
apply_background(s6)
add_header(s6, "Feature Highlight 03", "3D Algorithm Step-by-Step Voxel Visualizer", "Making time complexity, loop invariants, and data structures visually intuitive.")

tb6 = s6.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(4.8), Inches(4.6))
tf6 = tb6.text_frame
tf6.word_wrap = True

bullets_s6 = [
    ("Voxel Pillar Mapping:", "Array elements are rendered as 3D physical columns whose heights correspond to values."),
    ("Dynamic Color-Coded States:", "🟦 Default Blue  |  🟨 Comparing Amber  |  🟥 Active Swap Crimson  |  🟩 Sorted Emerald Green."),
    ("Full Playback Control:", "Step Forward, Step Backward, Play/Pause, Reset, and 1x / 2x / 4x speed toggles."),
    ("Theory & Invariant Panel:", "Explains the loop invariant (why the green zone grows from right to left) and compares worst-case O(n²) with best-case O(n)."),
    ("Deterministic Trace:", "Server calculates mathematical step traces deterministically before rendering.")
]
for t, d in bullets_s6:
    p = tf6.add_paragraph() if tf6.paragraphs[0].text else tf6.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

img_algo = os.path.join(IMAGES_DIR, "3d_algo_explained.png")
if os.path.exists(img_algo):
    s6.shapes.add_picture(img_algo, Inches(5.8), Inches(2.2), Inches(6.7), Inches(4.5))

# ==========================================
# SLIDE 7: SYSTEM 4 - 3D MATH & ASTROPHYSICS
# ==========================================
s7 = prs.slides.add_slide(blank_slide_layout)
apply_background(s7)
add_header(s7, "Feature Highlight 04", "3D Math Topology & Astrophysics Simulation", "Exploring multivariable calculus and orbital mechanics in photorealistic WebGL.")

# Math box
c_m = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.6), Inches(4.6))
c_m.fill.solid()
c_m.fill.fore_color.rgb = CARD_BG
c_m.line.color.rgb = BORDER_COLOR
tf_m = c_m.text_frame
tf_m.word_wrap = True
tf_m.margin_top = Inches(0.2)
tf_m.margin_left = Inches(0.3)
tf_m.margin_right = Inches(0.3)

p = tf_m.paragraphs[0]
p.text = "📐 3D Topological Math Surface Plotter"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf_m.add_paragraph()
p.text = "• Closed-form multivariable wave equation: z = 2.8·sin(1.1r)·e^(-0.07r)\n• Spectral elevation colormapping from ocean troughs to white crests.\n• Explains gradient vector fields (∇z) and stationary critical ring circles."
p.font.size = Pt(11)
p.font.color.rgb = MUTED
p.space_before = Pt(6)

img_math = os.path.join(IMAGES_DIR, "3d_math_explained.png")
if os.path.exists(img_math):
    s7.shapes.add_picture(img_math, Inches(0.9), Inches(3.8), Inches(5.4), Inches(2.8))

# Space box
c_s = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.2), Inches(5.6), Inches(4.6))
c_s.fill.solid()
c_s.fill.fore_color.rgb = CARD_BG
c_s.line.color.rgb = BORDER_COLOR
tf_s = c_s.text_frame
tf_s.word_wrap = True
tf_s.margin_top = Inches(0.2)
tf_s.margin_left = Inches(0.3)
tf_s.margin_right = Inches(0.3)

p = tf_s.paragraphs[0]
p.text = "🪐 AI-Generated 3D Solar System Simulation"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = INDIGO

p = tf_s.add_paragraph()
p.text = "• Gemini generates structured astrophysical scenes with real orbital radii.\n• Procedural planet textures, atmospheric Rayleigh scattering, and Saturn rings.\n• In-scene card explains Kepler's 3rd Law (T² ∝ a³) and orbital velocities."
p.font.size = Pt(11)
p.font.color.rgb = MUTED
p.space_before = Pt(6)

img_space = os.path.join(IMAGES_DIR, "3d_solar_system_explained.png")
if os.path.exists(img_space):
    s7.shapes.add_picture(img_space, Inches(6.9), Inches(3.8), Inches(5.4), Inches(2.8))

# ==========================================
# SLIDE 8: 13 TOOLS, 6 WIDGETS & VERIFICATION
# ==========================================
s8 = prs.slides.add_slide(blank_slide_layout)
apply_background(s8)
add_header(s8, "Verification & Impact", "13 Registered Tools, 6 Widgets & 100% Test Coverage", "Built with rock-solid architectural verification and multi-modal rigor.")

tb8 = s8.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(5.2), Inches(4.6))
tf8 = tb8.text_frame
tf8.word_wrap = True

bullets_s8 = [
    ("Complete Tool Suite (13 Tools):", "Physics playground, audio synth, algo visualizer, web sandbox, 3D math surface, Gemini 3D space scene, Python scratchpad, student analytics, code mission, unit test runner, simulation checker, Socratic hints, and progress store."),
    ("6 Generative UI Widgets:", "All bundled self-contained HTML/TypeScript widgets with zero client dependency issues."),
    ("Deterministic Trust Boundary:", "AI never grades its own homework. Real test executions determine progress."),
    ("Socratic Learning Loop:", "Students are guided to self-correction through real failure traces.")
]
for t, d in bullets_s8:
    p = tf8.add_paragraph() if tf8.paragraphs[0].text else tf8.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

img_report = os.path.join(IMAGES_DIR, "test_report_card.png")
if os.path.exists(img_report):
    s8.shapes.add_picture(img_report, Inches(6.2), Inches(2.2), Inches(6.3), Inches(4.5))

# ==========================================
# SLIDE 9: SUMMARY & VISION
# ==========================================
s9 = prs.slides.add_slide(blank_slide_layout)
apply_background(s9)
add_header(s9, "Conclusion & Horizon", "The Future of AI-Assisted Learning with Model Context Protocol", "Transforming LLMs from static text generators into interactive, multi-sensory physical laboratories.")

summary_cards = [
    ("Experiential Over Static", "Students learn 3-5x faster by manipulating physical parameters, listening to sound waves, and watching algorithms rather than reading walls of text.", CYAN, Inches(0.8), Inches(2.3)),
    ("Zero-Hallucination Integrity", "Rigid separation of concerns: Generative AI for creative visual modeling and Socratic questioning; deterministic engines for truth and grading.", EMERALD, Inches(4.8), Inches(2.3)),
    ("Extensible MCP Standard", "Built completely on open Model Context Protocol standards. Any MCP-compatible host (ChatGPT, Claude, IDEs) can immediately consume the tools and widgets.", INDIGO, Inches(8.8), Inches(2.3))
]

for title, desc, col, x, y in summary_cards:
    c = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.7), Inches(3.5))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.4)
    tf.margin_left = Inches(0.35)
    tf.margin_right = Inches(0.35)
    
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = col
    
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(14)

# Footer bar on Slide 9
tb_f = s9.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.6))
tf_f = tb_f.text_frame
p_f = tf_f.paragraphs[0]
p_f.alignment = PP_ALIGN.CENTER
p_f.text = "OmniLab MCP  •  The WebMCP Challenge  •  Open Source on GitHub: github.com/muhammad-hassaan-y2/webdev-MCP-Server"
p_f.font.size = Pt(12)
p_f.font.bold = True
p_f.font.color.rgb = CYAN

output_pptx = r"C:\Users\Hassaan\Downloads\mcp-tutor-py\OmniLab_WebMCP_Showcase.pptx"
prs.save(output_pptx)
print(f"SUCCESS: Generated presentation at {output_pptx}")
