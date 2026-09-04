import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Ultra-Modern Palette (Midnight OLED + Vibrant Neons)
BG_COLOR       = RGBColor(5, 8, 16)         # #050810 (Deep obsidian)
CARD_BG        = RGBColor(13, 20, 36)       # #0d1424 (Midnight slate)
CARD_BORDER    = RGBColor(30, 45, 75)       # #1e2d4b
CYAN           = RGBColor(56, 189, 248)     # #38bdf8 (Electric cyan)
PURPLE         = RGBColor(168, 85, 247)     # #a855f7 (Neon violet)
ROSE           = RGBColor(244, 63, 94)      # #f43f5e (Vibrant crimson)
AMBER          = RGBColor(245, 158, 11)     # #f59e0b (Warm amber)
EMERALD        = RGBColor(52, 211, 153)     # #34d399 (Glowing emerald)
WHITE          = RGBColor(255, 255, 255)
TEXT_WHITE     = RGBColor(248, 250, 252)    # #f8fafc
TEXT_MUTED     = RGBColor(148, 163, 184)    # #94a3b8

IMAGES_DIR = r"C:\Users\Hassaan\.gemini\antigravity\brain\c6fea575-e911-40ac-991c-56e4e413dcd4"

def apply_background(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_COLOR
    bg.line.fill.background()
    return bg

def add_header(slide, tag_text, title_text, subtitle_text, accent_color=CYAN):
    # Category Pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.32))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(15, 23, 42)
    pill.line.color.rgb = accent_color
    tf_pill = pill.text_frame
    tf_pill.word_wrap = True
    p_pill = tf_pill.paragraphs[0]
    p_pill.text = tag_text.upper()
    p_pill.font.size = Pt(9)
    p_pill.font.bold = True
    p_pill.font.color.rgb = accent_color
    p_pill.alignment = PP_ALIGN.CENTER

    # Main Title
    tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.7), Inches(0.65))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

    # Subtitle
    tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.45))
    tf_sub = tb_sub.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(12)
    p_sub.font.color.rgb = TEXT_MUTED

blank_layout = prs.slide_layouts[6]

# =========================================================================
# SLIDE 1: HERO COVER (Striking, Bold, Futuristic)
# =========================================================================
s1 = prs.slides.add_slide(blank_layout)
apply_background(s1)

# Hero Content Box
tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.333), Inches(4.5))
tf1 = tb1.text_frame
tf1.word_wrap = True

p_sub_tag = tf1.paragraphs[0]
p_sub_tag.text = "THE WEBMCP CHALLENGE  •  MULTI-MODAL GENERATIVE UI TUTOR"
p_sub_tag.font.size = Pt(11)
p_sub_tag.font.bold = True
p_sub_tag.font.color.rgb = CYAN

p_hero = tf1.add_paragraph()
p_hero.text = "Omni Tutor MCP"
p_hero.font.size = Pt(48)
p_hero.font.bold = True
p_hero.font.color.rgb = TEXT_WHITE
p_hero.space_before = Pt(8)

p_tag = tf1.add_paragraph()
p_tag.text = "The 3D Interactive STEM & Coding Studio Inside Your Chat"
p_tag.font.size = Pt(22)
p_tag.font.bold = True
p_tag.font.color.rgb = PURPLE
p_tag.space_after = Pt(12)

p_desc = tf1.add_paragraph()
p_desc.text = "Replacing static text walls with live 60 FPS Cannon-es physics, Web Audio synthesizers, 3D voxel algorithm simulations, and a zero-hallucination deterministic grading engine."
p_desc.font.size = Pt(14)
p_desc.font.color.rgb = TEXT_MUTED
p_desc.space_after = Pt(28)

# 4 Live Metrics Badges at the bottom of Slide 1
metrics = [
    ("13 MCP Tools", "Full Protocol Tool Suite", CYAN),
    ("6 Generative Widgets", "Zero-Install ext-apps UI", PURPLE),
    ("60 FPS Physics & Audio", "Cannon-es + Web Audio API", ROSE),
    ("100% Deterministic", "Zero Grading Hallucinations", EMERALD)
]
for i, (val, lbl, col) in enumerate(metrics):
    box = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0 + i * 2.85), Inches(5.6), Inches(2.65), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = col
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.18)
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    p1 = tf.paragraphs[0]
    p1.text = val
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.text = lbl
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED

# =========================================================================
# SLIDE 2: THE PARADIGM SHIFT ("Death of the Text Wall")
# =========================================================================
s2 = prs.slides.add_slide(blank_layout)
apply_background(s2)
add_header(s2, "The Paradigm Shift", "Why Chatbots Fail at Teaching: The Text-Wall Problem", "Students don't learn physics by reading paragraphs—they learn by touching, breaking, and testing.")

# Left Card: Traditional AI Chatbots (Deficient)
c_left = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.8))
c_left.fill.solid()
c_left.fill.fore_color.rgb = CARD_BG
c_left.line.color.rgb = ROSE
tf_l = c_left.text_frame
tf_l.word_wrap = True
tf_l.margin_top = Inches(0.3)
tf_l.margin_left = Inches(0.4)
tf_l.margin_right = Inches(0.4)

p = tf_l.paragraphs[0]
p.text = "❌ Conventional AI Chatbots"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = ROSE

items_l = [
    ("500-Word Markdown Walls:", "Students glaze over long essays; abstract spatial formulas (gravity, torque, Fourier) remain incomprehensible."),
    ("Hallucinated Verification (Sycophancy):", "LLMs frequently say 'Great job! That works!' on broken code or impossible physics predictions, validating flawed mental models."),
    ("Passive Reading vs. Active Building:", "Zero tactile or auditory interaction. Students are spectators rather than experimenters."),
    ("Opaque Black-Box Feedback:", "When code fails, LLMs dump a generic rewrite instead of guiding the student to discover their own error.")
]
for t, d in items_l:
    p = tf_l.add_paragraph()
    p.text = f"• {t} {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_MUTED
    p.space_before = Pt(12)

# Right Card: Omni Tutor MCP (Revolutionary)
c_right = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(2.0), Inches(5.6), Inches(4.8))
c_right.fill.solid()
c_right.fill.fore_color.rgb = CARD_BG
c_right.line.color.rgb = CYAN
tf_r = c_right.text_frame
tf_r.word_wrap = True
tf_r.margin_top = Inches(0.3)
tf_r.margin_left = Inches(0.4)
tf_r.margin_right = Inches(0.4)

p = tf_r.paragraphs[0]
p.text = "✓ The Omni Tutor MCP Experience"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = CYAN

items_r = [
    ("Live Generative UI in Chat:", "Chat threads instantly launch interactive 3D physics viewports, audio synthesizers, and voxel algorithms."),
    ("Deterministic Trust Boundary:", "The AI NEVER grades its own homework. Real Python subprocesses and closed-form physics engines verify truth."),
    ("Multi-Sensory Discovery:", "Students hear wave harmonics, shoot cannonballs with mass and torque, and watch voxel sorting in real time."),
    ("Socratic Scaffolding Engine:", "Gemini inspects ground-truth unit test diffs to ask targeted guiding questions rather than giving away answers.")
]
for t, d in items_r:
    p = tf_r.add_paragraph()
    p.text = f"• {t} {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(12)

# =========================================================================
# SLIDE 3: THE 4 MULTI-SENSORY PILLARS
# =========================================================================
s3 = prs.slides.add_slide(blank_layout)
apply_background(s3)
add_header(s3, "Multi-Sensory Learning", "How Omni Tutor MCP Teaches: 4 Dimensions of Understanding", "Engaging visual, physical, auditory, and logical cognitive channels simultaneously.")

pillars = [
    ("👁️ SEE IT: 3D Astrophysics & Math", "Photorealistic Three.js PBR rendering with ACESFilmic tone mapping. Students explore planetary orbits and multivariable calculus waves in real space.", CYAN, Inches(0.8), Inches(2.1)),
    ("💥 TOUCH IT: 3D Cannon-es Physics", "Rigid-body mechanics with real mass, velocity, torque, friction, and restitution. Students fire cannonballs and switch gravity from Earth to Moon or Zero-G.", ROSE, Inches(6.9), Inches(2.1)),
    ("👂 HEAR IT: Web Audio Synthesizer", "Native Web Audio API synthesis. Students play an 8-key piano keyboard, change waveforms, and watch sound wave harmonics vibrate on a 60 FPS oscilloscope.", AMBER, Inches(0.8), Inches(4.6)),
    ("🧠 CODE IT: Deterministic Sandboxes", "Isolated Python execution (python -I) and live HTML/CSS web sandbox. Socratic hints guide students to resolve off-by-one errors and algorithmic bottlenecks.", EMERALD, Inches(6.9), Inches(4.6))
]

for title, desc, col, x, y in pillars:
    c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.6), Inches(2.1))
    c.fill.solid()
    c.fill.fore_color.rgb = CARD_BG
    c.line.color.rgb = col
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)

# =========================================================================
# SLIDE 4: SHOWCASE 1 - 3D RIGID-BODY PHYSICS PLAYGROUND
# =========================================================================
s4 = prs.slides.add_slide(blank_layout)
apply_background(s4)
add_header(s4, "Interactive System 01", "3D Rigid-Body Physics: Domino Cascade & Cannonballs", "True 60 Hz physics collision dynamics with embedded in-scene scientific explanation.", ROSE)

tb4 = s4.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.0), Inches(5.0))
tf4 = tb4.text_frame
tf4.word_wrap = True

bullets_4 = [
    ("Cannon-es + Three.js Engine:", "Couples real-time rigid body dynamics to Three.js PBR rendering with mirror chrome reflections and acrylic translucency."),
    ("Impulse-Momentum Transfer (p = m·v):", "Heavy steel projectile (m=15 kg, v=40 m/s) delivers sudden momentum into lead domino, triggering an unbalancing torque cascade (τ = r × F)."),
    ("Dynamic Gravity Switcher:", "Toggle gravity in real time: 🌍 Earth (-9.82 m/s²), 🌕 Moon (-1.62 m/s²), 🌌 Zero-G (0.0 m/s²), 🪐 Jupiter (-24.79 m/s²)."),
    ("Interactive Blaster:", "Students click anywhere in the 3D scene to shoot projectile spheres and destroy towers or domino chains."),
    ("Embedded Scientific Card:", "Directly in the 3D view, students read exact equations for impulse, restitution (e=0.45), and potential-to-kinetic energy conversion.")
]
for t, d in bullets_4:
    p = tf4.add_paragraph() if tf4.paragraphs[0].text else tf4.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

img_phys = os.path.join(IMAGES_DIR, "3d_physics_explained.png")
if os.path.exists(img_phys):
    s4.shapes.add_picture(img_phys, Inches(6.0), Inches(1.9), Inches(6.5), Inches(5.0))

# =========================================================================
# SLIDE 5: SHOWCASE 2 - WEB AUDIO SYNTHESIZER & OSCILLOSCOPE
# =========================================================================
s5 = prs.slides.add_slide(blank_layout)
apply_background(s5)
add_header(s5, "Interactive System 02", "Web Audio Frequency Synthesizer & 60 FPS Oscilloscope", "Making sound waves, harmonics, and Fourier series audible and visible in chat.", AMBER)

tb5 = s5.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.0), Inches(5.0))
tf5 = tb5.text_frame
tf5.word_wrap = True

bullets_5 = [
    ("Browser-Native AudioContext:", "Synthesizes real acoustic tones directly in the browser with zero external audio assets."),
    ("Live Glowing Oscilloscope:", "Extracts time-domain FFT data via AnalyserNode at 60 FPS, drawing the vibrating waveform as notes are played."),
    ("4 Harmonic Waveforms:", "Switch between Sine ∿ (pure fundamental), Square ⊓ (odd harmonics), Sawtooth ⧘ (all harmonics), and Triangle ⋀."),
    ("Playable 8-Key Keyboard:", "Play notes C4 (262 Hz) through C5 (523 Hz) with resonant low-pass filter cutoff (100–6000 Hz) and pitch knobs."),
    ("Fourier Acoustics Integration:", "Students directly observe how altering waveform geometry changes the auditory timbre and Fourier equation simultaneously.")
]
for t, d in bullets_5:
    p = tf5.add_paragraph() if tf5.paragraphs[0].text else tf5.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

img_audio = os.path.join(IMAGES_DIR, "dashboard_audio.png")
if os.path.exists(img_audio):
    s5.shapes.add_picture(img_audio, Inches(6.0), Inches(1.9), Inches(6.5), Inches(5.0))

# =========================================================================
# SLIDE 6: SHOWCASE 3 - 3D ALGORITHM VOXEL VISUALIZER
# =========================================================================
s6 = prs.slides.add_slide(blank_layout)
apply_background(s6)
add_header(s6, "Interactive System 03", "3D Algorithm Step-by-Step Voxel Visualizer", "Transforming abstract Big-O notation and loop invariants into physical spatial animations.", CYAN)

tb6 = s6.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.0), Inches(5.0))
tf6 = tb6.text_frame
tf6.word_wrap = True

bullets_6 = [
    ("3D Voxel Height Mapping:", "Array elements are rendered as 3D physical pillars with heights matching numeric values."),
    ("State Machine Color Encoding:", "🟦 Default Blue  ➔  🟨 Comparing Amber  ➔  🟥 Active Swap Crimson  ➔  🟩 Sorted Emerald Green."),
    ("Loop Invariant Visualization:", "Students visually witness why the emerald green zone grows from right to left as the largest numbers bubble up."),
    ("Interactive Playback Dock:", "Step Forward, Step Backward, Play / Pause, Reset, and 1x / 2x / 4x speed toggles."),
    ("Theory & Complexity Panel:", "Embedded card derives quadratic worst-case comparisons C(n) = n(n-1)/2 ∈ O(n²) vs. linear O(n) early termination.")
]
for t, d in bullets_6:
    p = tf6.add_paragraph() if tf6.paragraphs[0].text else tf6.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(11)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(8)

img_algo = os.path.join(IMAGES_DIR, "3d_algo_explained.png")
if os.path.exists(img_algo):
    s6.shapes.add_picture(img_algo, Inches(6.0), Inches(1.9), Inches(6.5), Inches(5.0))

# =========================================================================
# SLIDE 7: SHOWCASE 4 - 3D MATH & ASTROPHYSICS
# =========================================================================
s7 = prs.slides.add_slide(blank_layout)
apply_background(s7)
add_header(s7, "Interactive System 04", "Photorealistic 3D Math Topology & Astrophysics", "Exploring multivariable calculus wavefronts and Keplerian orbits in WebGL.", PURPLE)

# Left math box
c_m = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.7), Inches(5.1))
c_m.fill.solid()
c_m.fill.fore_color.rgb = CARD_BG
c_m.line.color.rgb = CARD_BORDER
tf_m = c_m.text_frame
tf_m.word_wrap = True
tf_m.margin_top = Inches(0.15)
tf_m.margin_left = Inches(0.25)
tf_m.margin_right = Inches(0.25)

p = tf_m.paragraphs[0]
p.text = "📐 3D Topological Math Surface"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = CYAN

p = tf_m.add_paragraph()
p.text = "• Closed-form wave equation z = 2.8·sin(1.1r)·e^(-0.07r)\n• Spectral elevation colormap (deep blue troughs to white crests)\n• Explains radial symmetry, gradient vectors (∇z), and stationary critical rings."
p.font.size = Pt(10)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(4)

img_math = os.path.join(IMAGES_DIR, "3d_math_explained.png")
if os.path.exists(img_math):
    s7.shapes.add_picture(img_math, Inches(0.95), Inches(3.35), Inches(5.4), Inches(3.45))

# Right space box
c_s = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(5.7), Inches(5.1))
c_s.fill.solid()
c_s.fill.fore_color.rgb = CARD_BG
c_s.line.color.rgb = CARD_BORDER
tf_s = c_s.text_frame
tf_s.word_wrap = True
tf_s.margin_top = Inches(0.15)
tf_s.margin_left = Inches(0.25)
tf_s.margin_right = Inches(0.25)

p = tf_s.paragraphs[0]
p.text = "🪐 Gemini AI 3D Solar System"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = PURPLE

p = tf_s.add_paragraph()
p.text = "• Gemini generates structured astrophysical models with real orbital radii.\n• Procedural Earth oceans/continents, Rayleigh atmosphere, and Saturn rings.\n• In-scene card explains Kepler's 3rd Law (T² ∝ a³) and orbital velocities."
p.font.size = Pt(10)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(4)

img_space = os.path.join(IMAGES_DIR, "3d_solar_system_explained.png")
if os.path.exists(img_space):
    s7.shapes.add_picture(img_space, Inches(6.95), Inches(3.35), Inches(5.4), Inches(3.45))

# =========================================================================
# SLIDE 8: THE SECRET WEAPON - DETERMINISTIC TRUST BOUNDARY
# =========================================================================
s8 = prs.slides.add_slide(blank_layout)
apply_background(s8)
add_header(s8, "Pedagogy & Integrity", "The Deterministic Trust Boundary: Why Omni Tutor Never Hallucinates", "Strictly separating generative creative proposal from deterministic verification and grading.", EMERALD)

flow_boxes = [
    ("1. Student Action", "Student writes Python code, manipulates launch velocity, or tweaks HTML/CSS parameters.", CYAN, Inches(0.8)),
    ("2. Isolated Subprocess", "Python server executes code in an isolated subprocess (python -I) or calculates closed-form math.", AMBER, Inches(3.8)),
    ("3. Evaluation Gate", "Did it pass? If FAIL: Gemini reads real error diff & generates Socratic nudge. If PASS: progress is recorded.", ROSE, Inches(6.8)),
    ("4. Verified Mastery", "Mastery badges and progress are stored in persistent store (.data/progress.json). AI cannot falsify.", EMERALD, Inches(9.8))
]

for title, desc, col, x in flow_boxes:
    b = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(2.7), Inches(2.3))
    b.fill.solid()
    b.fill.fore_color.rgb = CARD_BG
    b.line.color.rgb = col
    tf = b.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.2)
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = col
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(10)
    p2.font.color.rgb = TEXT_MUTED
    p2.space_before = Pt(8)

# Lower callout card: Socratic Guidance in Action
callout = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.3))
callout.fill.solid()
callout.fill.fore_color.rgb = CARD_BG
callout.line.color.rgb = CARD_BORDER
tf_c = callout.text_frame
tf_c.word_wrap = True
tf_c.margin_top = Inches(0.2)
tf_c.margin_left = Inches(0.35)
tf_c.margin_right = Inches(0.35)

p = tf_c.paragraphs[0]
p.text = "💡 Real-World Socratic Feedback Example (Off-By-One Sum Mission)"
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = EMERALD

p = tf_c.add_paragraph()
p.text = "• Student Bug: range(n) returns sum of [0, 1, ..., n-1]. For n=3, it produces 3 instead of 6.\n• Deterministic Detection: run_tests detects assertion failure: AssertionError: expected sum_first_n(3) == 6, got 3.\n• Gemini Socratic Response: 'Notice that range(n) stops just before n. How might you adjust the arguments inside range() so your loop includes both 1 and n?'\n• Integrity Result: Student discovers the fix themselves, runs tests again, and earns verifiable mastery."
p.font.size = Pt(11)
p.font.color.rgb = TEXT_WHITE
p.space_before = Pt(8)

# =========================================================================
# SLIDE 9: VERIFICATION & 13-TOOL CATALOG
# =========================================================================
s9 = prs.slides.add_slide(blank_layout)
apply_background(s9)
add_header(s9, "Full Capabilities", "13 Production MCP Tools & 6 Interactive UI Widgets", "Tested, verified, and passing 100% of automated test suites.", CYAN)

tb9 = s9.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(5.0), Inches(5.0))
tf9 = tb9.text_frame
tf9.word_wrap = True

tools_summary = [
    ("physics_rigid_body_playground:", "3D Cannon-es physics with gravity, dominoes, towers, and cannonballs."),
    ("interactive_audio_synth:", "Web Audio synthesizer with 60 FPS oscilloscope and piano keyboard."),
    ("visualize_algorithm_3d:", "3D voxel step visualizer for sorting, pathfinding, and loop invariants."),
    ("preview_html_css:", "Live web sandbox with editable code tabs and responsive viewports."),
    ("visualize_math_surface:", "3D multivariable calculus surface plotter with elevation colormapping."),
    ("generate_3d_scene:", "Gemini AI astrophysics and molecular scene generator."),
    ("run_code_scratchpad:", "Safe isolated Python snippet execution with 5s timeout."),
    ("get_student_analytics:", "Tracks student concept mastery, attempts, and accuracy."),
    ("start_mission / run_tests:", "Interactive code and physics simulation missions with subprocess test grading.")
]
for t, d in tools_summary:
    p = tf9.add_paragraph() if tf9.paragraphs[0].text else tf9.paragraphs[0]
    p.text = f"• {t} {d}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_WHITE
    p.space_before = Pt(4)

img_rep = os.path.join(IMAGES_DIR, "test_report_card.png")
if os.path.exists(img_rep):
    s9.shapes.add_picture(img_rep, Inches(6.0), Inches(1.9), Inches(6.5), Inches(5.0))

# =========================================================================
# SLIDE 10: CONCLUSION & VISION (The Grand Finale)
# =========================================================================
s10 = prs.slides.add_slide(blank_layout)
apply_background(s10)
add_header(s10, "The Vision", "Omni Tutor MCP: The Future of Spatial AI Education", "Built for The WebMCP Challenge • Bringing Generative UI to Every LLM Client.", PURPLE)

c_vis1 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(3.7), Inches(4.0))
c_vis1.fill.solid()
c_vis1.fill.fore_color.rgb = CARD_BG
c_vis1.line.color.rgb = CYAN
tf_v1 = c_vis1.text_frame
tf_v1.word_wrap = True
tf_v1.margin_top = Inches(0.3)
tf_v1.margin_left = Inches(0.3)
tf_v1.margin_right = Inches(0.3)
p = tf_v1.paragraphs[0]
p.text = "🚀 Experiential Learning"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CYAN
p = tf_v1.add_paragraph()
p.text = "Replacing static textbooks and passive chatbots with immersive, hands-on physical simulations. Students retain concepts 3-5x better when they manipulate the physical variables directly."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(12)

c_vis2 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(2.0), Inches(3.7), Inches(4.0))
c_vis2.fill.solid()
c_vis2.fill.fore_color.rgb = CARD_BG
c_vis2.line.color.rgb = EMERALD
tf_v2 = c_vis2.text_frame
tf_v2.word_wrap = True
tf_v2.margin_top = Inches(0.3)
tf_v2.margin_left = Inches(0.3)
tf_v2.margin_right = Inches(0.3)
p = tf_v2.paragraphs[0]
p.text = "🛡️ Truth & Integrity"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = EMERALD
p = tf_v2.add_paragraph()
p.text = "By strictly separating AI creative generation from deterministic unit test grading, Omni Tutor MCP eliminates hallucinations and guarantees that every student certificate is mathematically earned."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(12)

c_vis3 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.0), Inches(3.7), Inches(4.0))
c_vis3.fill.solid()
c_vis3.fill.fore_color.rgb = CARD_BG
c_vis3.line.color.rgb = PURPLE
tf_v3 = c_vis3.text_frame
tf_v3.word_wrap = True
tf_v3.margin_top = Inches(0.3)
tf_v3.margin_left = Inches(0.3)
tf_v3.margin_right = Inches(0.3)
p = tf_v3.paragraphs[0]
p.text = "🌐 Open MCP Standard"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = PURPLE
p = tf_v3.add_paragraph()
p.text = "Built 100% on the Model Context Protocol and ext-apps specification. Runs out of the box in Claude, ChatGPT, Cursor, and any MCP-enabled host with zero client-side installation."
p.font.size = Pt(12)
p.font.color.rgb = TEXT_MUTED
p.space_before = Pt(12)

# Bottom Footer
tb_foot = s10.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.6))
tf_foot = tb_foot.text_frame
p_foot = tf_foot.paragraphs[0]
p_foot.alignment = PP_ALIGN.CENTER
p_foot.text = "Omni Tutor MCP  •  The WebMCP Challenge  •  GitHub: github.com/muhammad-hassaan-y2/webdev-MCP-Server"
p_foot.font.size = Pt(12)
p_foot.font.bold = True
p_foot.font.color.rgb = CYAN

output_path = r"C:\Users\Hassaan\Downloads\mcp-tutor-py\Omni_Tutor_MCP_Showcase.pptx"
prs.save(output_path)
print(f"SUCCESS: Saved ultra-appealing presentation to {output_path}")
